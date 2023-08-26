from pptx import Presentation
import os
import doc2docx
from docx2python import docx2python
import docx2txt
from pdf2image import convert_from_path


def convert_formats(fil):
    """
    Function converts input files 
    into allowed image formats
    """
    folders = os.listdir('./data/')
    for folder in folders:

        images = os.listdir(f'./data/{folder}/images')

        for image in images:
            filename, ext = image.split('.')
            extl = ext.lower()

            if extl == 'png' or extl == 'jpg' or extl == 'jpeg':
                os.system(f'cp data/{folder}/images/{filename}.{ext} img/{folder}/{filename}.{ext}')

            elif extl.startswith('doc'):
                if extl == 'doc':
                    doc2docx.convert(f'./data/{folder}/images/{filename}.{ext}', f'./data/{folder}/images/{filename}.docx')
                doc_res = docx2python(f'./data/{folder}/images/{filename}.docx')
                for key in doc_res.images:
                    b = doc_res.images[key]
                    ext = key.split('.')[-1]
                    with open(f'./img/{folder}/{filename}.{ext}', 'wb') as file:
                        file.write(b)
                    break
            elif extl == 'pptx':
                pres = Presentation(f'./data/{folder}/images/{filename}.{ext}')
                slide = pres.slides[0]
                shape = slide.shapes[0]
                image = shape.image
                blob = image.blob
                ext = image.ext
                with open(f'./img/{filename}.{ext}', 'wb') as file:
                    file.write(blob)
            elif extl == 'pdf':
                pages = convert_from_path(f'./data/{folder}/images/{filename}.pdf', 500)
                page = pages[0]
                page.save(f'./img/{folder}/{filename}.jpg', 'JPEG')
